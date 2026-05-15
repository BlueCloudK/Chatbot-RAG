import os
import re
import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from langchain_text_splitters import RecursiveCharacterTextSplitter


class DocumentProcessor:
    def __init__(self, chunk_size=1000, chunk_overlap=200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            length_function=len,
            separators=["\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " ", ""]
        )

    def clean_text(self, text):
        """Làm sạch text: bỏ whitespace thừa, ký tự lạ, normalize unicode"""
        # Bỏ ký tự điều khiển
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
        # Bỏ khoảng trắng thừa
        text = re.sub(r'[ \t]+', ' ', text)
        # Bỏ dòng trống thừa (giữ tối đa 2 dòng trống)
        text = re.sub(r'\n{3,}', '\n\n', text)
        # Bỏ khoảng trắng đầu/cuối mỗi dòng
        lines = [line.strip() for line in text.split('\n')]
        text = '\n'.join(lines)
        return text.strip()

    def extract_text_from_pdf(self, file_path):
        text = ""
        with open(file_path, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        return text

    def extract_text_from_docx(self, file_path):
        doc = DocxDocument(file_path)
        parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        return "\n".join(parts)

    def extract_text_from_pptx(self, file_path):
        prs = Presentation(file_path)
        text = ""
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                text += f"\n--- Slide {slide_num} ---\n"
                text += "\n".join(slide_text) + "\n"
        return text

    def process_file(self, file_path):
        """Đọc file, clean text, cắt thành chunks"""
        ext = os.path.splitext(file_path)[1].lower()
        if ext == '.pdf':
            text = self.extract_text_from_pdf(file_path)
        elif ext == '.docx':
            text = self.extract_text_from_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            text = self.extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

        # Clean text trước khi chunk
        text = self.clean_text(text)
        
        if not text or len(text) < 50:
            print(f"  Warning: Very little text extracted ({len(text)} chars)")
            return []

        chunks = self.text_splitter.split_text(text)
        
        print(f"  Extracted {len(text)} chars → {len(chunks)} chunks (size={self.chunk_size}, overlap={self.chunk_overlap})")
        return chunks
